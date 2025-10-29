"""
FastAPI PPTX Microservice
- /           : root ping
- /health     : health check
- /template/analyze : analyze PPTX template -> style profile
- /deck/generate    : generate PPTX from DeckSpec + templateUrl (returns .pptx file)
- /deck/patch       : apply PatchOps to an existing deck (returns .pptx file)
"""

from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Dict, Any, Optional
from pathlib import Path
import os
import shutil
import tempfile
import json
import urllib.request

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# --- App & CORS --------------------------------------------------------------

app = FastAPI(title="PPTX Service", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Enable slash tolerance
app.router.redirect_slashes = True

# Temp working dir
TEMP_DIR = Path(tempfile.gettempdir()) / "pptx-service"
TEMP_DIR.mkdir(exist_ok=True)

# --- Startup logs (widoczne w Render -> Logs) --------------------------------
@app.on_event("startup")
async def _startup_log():
    try:
        paths = [getattr(r, "path", str(r)) for r in app.router.routes]
        print("🚀 Starting PPTX Service...")
        print("✅ Routes registered:", paths)
        print("📂 TEMP_DIR:", str(TEMP_DIR))
    except Exception as e:
        print("Startup log error:", e)

# --- Models ------------------------------------------------------------------

class DeckSpec(BaseModel):
    title: Optional[str] = ""
    slides: List[Dict[str, Any]] = []
    meta: Optional[Dict[str, str]] = {"locale": "en-US", "timezone": "Europe/Warsaw"}

class PatchOps(BaseModel):
    ops: List[Dict[str, Any]]

class GenerateRequest(BaseModel):
    templateUrl: str           # public URL do .pptx template
    deckSpec: Dict[str, Any]   # elastyczny DeckSpec (LLM output)

class PatchRequest(BaseModel):
    pptxUrl: str               # publiczny URL do istniejącego .pptx
    patchOps: PatchOps
    currentSpec: Dict[str, Any]

# --- Ping / Health -----------------------------------------------------------

@app.get("/")
def root():
    return {"status": "ok", "service": "pptx-service"}

@app.get("/health")
def health():
    return {"status": "healthy", "service": "pptx-service", "version": "1.0.0"}

# --- Analyze template --------------------------------------------------------

@app.post("/template/analyze")
async def analyze_template(file: UploadFile = File(...)):
    """
    Upload .pptx -> extract slide size, layouts, placeholders, basic theme colors
    Returns a JSON "style profile" used later by the LLM Deck Planner.
    """
    temp_path = TEMP_DIR / f"template_{os.urandom(8).hex()}.pptx"
    try:
        # Save uploaded file
        with temp_path.open("wb") as f:
            shutil.copyfileobj(file.file, f)

        prs = Presentation(temp_path)

        style_profile: Dict[str, Any] = {
            "slideSize": {"width": prs.slide_width, "height": prs.slide_height},
            "layouts": [],
            "themeColors": [],
            "fonts": [],
        }

        # Layouts & placeholders
        for idx, layout in enumerate(prs.slide_layouts):
            layout_info = {
                "name": getattr(layout, "name", f"Layout {idx}"),
                "index": idx,
                "placeholders": [],
            }
            for ph in layout.placeholders:
                # typ placeholdera może być Enum (serializujemy do str)
                ph_type = getattr(getattr(ph, "placeholder_format", None), "type", None)
                layout_info["placeholders"].append({
                    "idx": getattr(getattr(ph, "placeholder_format", None), "idx", None),
                    "type": str(ph_type),
                    "name": getattr(ph, "name", None),
                })
            style_profile["layouts"].append(layout_info)

        # Theme colors (best-effort)
        try:
            theme = prs.part.theme_part.theme
            for c in theme.color_scheme:
                style_profile["themeColors"].append({
                    "name": getattr(c, "name", None),
                    "rgb": str(getattr(c, "rgb", "")),
                })
        except Exception:
            pass

        return style_profile

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            temp_path.unlink(missing_ok=True)
        except Exception:
            pass

# --- Generate deck -----------------------------------------------------------

@app.post("/deck/generate")
async def generate_deck(request: GenerateRequest):
    """
    Build a .pptx from the provided templateUrl and deckSpec.
    Returns the generated PPTX as a FileResponse (download).
    """
    temp_template = TEMP_DIR / f"template_{os.urandom(8).hex()}.pptx"
    output_path = TEMP_DIR / f"deck_{os.urandom(8).hex()}.pptx"

    try:
        # Download template
        urllib.request.urlretrieve(request.templateUrl, temp_template)

        prs = Presentation(temp_template)

        # Remove all existing slides from template, keep masters/theme
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        # Build slides from spec
        slides = request.deckSpec.get("slides", [])
        for slide_spec in slides:
            layout_name = slide_spec.get("layoutName") or slide_spec.get("layout") or "Title Slide"

            # find layout by name, fallback to first
            layout = None
            for l in prs.slide_layouts:
                if getattr(l, "name", "") == layout_name:
                    layout = l
                    break
            if not layout:
                layout = prs.slide_layouts[0]

            slide = prs.slides.add_slide(layout)

            # title (best-effort by shape name)
            title_txt = slide_spec.get("title", "")
            if title_txt:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and "title" in getattr(shape, "name", "").lower():
                        shape.text = title_txt
                        break

            # bullets (best-effort by "content" name)
            bullets = slide_spec.get("bullets", [])
            if bullets:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and "content" in getattr(shape, "name", "").lower():
                        tf = shape.text_frame
                        tf.clear()
                        for b in bullets:
                            p = tf.add_paragraph()
                            p.text = str(b)
                            p.level = 0
                        break

            # charts (minimal: bar/column/line/pie)
            if "chart" in slide_spec:
                add_chart_to_slide(slide, slide_spec["chart"])

        prs.save(output_path)

        # Return the file directly (frontend/Edge Function może przesłać dalej lub zapisać w Storage)
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="generated.pptx",
            headers={
                "X-Spec-JSON": json.dumps(request.deckSpec),  # pomocniczo do debug
                "X-Thumbnails": json.dumps([]),
            },
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            temp_template.unlink(missing_ok=True)
        except Exception:
            pass
        # NIE usuwamy output_path – musi zostać do czasu wysłania pliku

# --- Patch deck --------------------------------------------------------------

@app.post("/deck/patch")
async def patch_deck(request: PatchRequest):
    """
    Apply simple PatchOps to an existing deck and return a new .pptx file.
    Supported ops (MVP):
      - replace_text { slideIndex, placeholder, newText }
      - add_slide    { afterIndex?, layout, placeholders{ name:text } }
      - update_chart_series (stub)
    """
    temp_deck = TEMP_DIR / f"deck_{os.urandom(8).hex()}.pptx"
    output_path = TEMP_DIR / f"deck_patched_{os.urandom(8).hex()}.pptx"

    try:
        urllib.request.urlretrieve(request.pptxUrl, temp_deck)
        prs = Presentation(temp_deck)

        updated_spec = dict(request.currentSpec)

        for op in request.patchOps.ops:
            op_type = op.get("type")

            if op_type == "replace_text":
                slide_idx = int(op.get("slideIndex", 0))
                placeholder = op.get("placeholder")
                new_text = op.get("newText", "")
                if 0 <= slide_idx < len(prs.slides):
                    slide = prs.slides[slide_idx]
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False) and getattr(shape, "name", "") == placeholder:
                            shape.text = new_text
                            break

            elif op_type == "update_chart_series":
                # TODO: implement real chart data update if needed
                pass

            elif op_type == "add_slide":
                layout_name = op.get("layout", "Title and Content")
                layout = None
                for l in prs.slide_layouts:
                    if getattr(l, "name", "") == layout_name:
                        layout = l
                        break
                if not layout:
                    layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                placeholders = op.get("placeholders", {})
                for ph_name, text in placeholders.items():
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False) and getattr(shape, "name", "") == ph_name:
                            shape.text = str(text)
                            break

        prs.save(output_path)

        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="patched.pptx",
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            temp_deck.unlink(missing_ok=True)
        except Exception:
            pass

# --- Helpers -----------------------------------------------------------------

def add_chart_to_slide(slide, chart_spec: Dict[str, Any]):
    """
    chart_spec: { "type":"bar|column|line|pie", "title":"...", "data":[["Cat", value], ...] }
    """
    chart_type = chart_spec.get("type", "column")
    title = chart_spec.get("title", "")
    data = chart_spec.get("data", [])

    # Position & size
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)

    chart_data = CategoryChartData()
    chart_data.categories = [row[0] for row in data]
    chart_data.add_series("Series 1", [row[1] for row in data])

    chart_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_shape = slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, chart_data)
    chart = chart_shape.chart
    chart.has_legend = True
    if title:
        chart.chart_title.text_frame.text = title

# --- Dev run (nieużywane na Render, ale przydatne lokalnie) ------------------

if __name__ == "__main__":
    import uvicorn
    print("🔧 Running dev server on http://0.0.0.0:8000")
    uvicorn.run(app, host="0.0.0.0", port=8000)
